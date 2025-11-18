# step1_extraction.py (Updated with Manifest Creation)

import os
import json
from pathlib import Path
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from typing import Dict, Any

class DocumentExtractor:
    """Extract text from PDF, DOCX, PPTX, TXT, MD files"""
    
    def __init__(self, source_dir: str = "source_documents", 
                 output_dir: str = "output/extracted_texts"):
        self.source_dir = Path(source_dir)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True, parents=True)
        
        print(f"📂 Source directory: {self.source_dir}")
        print(f"📂 Output directory: {self.output_dir}\n")
    
    def extract_from_pdf(self, pdf_path: str) -> str:
        text = ""
        try:
            reader = PdfReader(pdf_path)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
        except Exception as e:
            print(f"      ❌ Error reading PDF {Path(pdf_path).name}: {e}")
        return text.strip()
    
    def extract_from_docx(self, docx_path: str) -> str:
        text = ""
        try:
            doc = DocxDocument(docx_path)
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        except Exception as e:
            print(f"      ❌ Error reading DOCX {Path(docx_path).name}: {e}")
        return text.strip()
    
    def extract_from_pptx(self, pptx_path: str) -> str:
        text = ""
        try:
            prs = Presentation(pptx_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text += f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text) + "\n"
        except Exception as e:
            print(f"      ❌ Error reading PPTX {Path(pptx_path).name}: {e}")
        return text.strip()
    
    def extract_from_text(self, text_path: str) -> str:
        text = ""
        try:
            with open(text_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except Exception as e:
            print(f"      ❌ Error reading text file {Path(text_path).name}: {e}")
        return text.strip()

    def extract_single_file(self, file_path: Path) -> Dict[str, Any]:
        ext = file_path.suffix.lower()
        extractors = {'.pdf': self.extract_from_pdf, '.docx': self.extract_from_docx,
                      '.pptx': self.extract_from_pptx, '.txt': self.extract_from_text,
                      '.md': self.extract_from_text}
        
        if ext not in extractors:
            return {'success': False, 'filename': file_path.name, 'error': 'Unsupported format'}
        
        text = extractors[ext](str(file_path))
        
        if not text:
            return {'success': False, 'filename': file_path.name, 'error': 'No text extracted'}
        
        return {'success': True, 'text': text, 'filename': file_path.name}

    def save_extracted_text(self, text: str, original_filename: str):
        output_filename = Path(original_filename).stem + ".txt"
        output_path = self.output_dir / output_filename
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return output_path

    def extract_all_documents(self):
        supported_extensions = {'.pdf', '.docx', '.pptx', '.txt', '.md'}
        all_files = [p for p in self.source_dir.rglob("*") if p.suffix.lower() in supported_extensions]
        
        if not all_files:
            print(f"⚠ No supported files found in {self.source_dir}")
            return

        print(f"📚 Found {len(all_files)} files to process\n")
        
        extraction_manifest = {}
        successful, failed = 0, 0
        
        for idx, file_path in enumerate(all_files, 1):
            print(f"[{idx}/{len(all_files)}] Processing: {file_path.name}")
            
            result = self.extract_single_file(file_path)
            
            if result['success']:
                output_path = self.save_extracted_text(result['text'], result['filename'])
                # Map the new .txt filename to the original source filename
                extraction_manifest[output_path.name] = result['filename']
                successful += 1
                print(f"      ✅ Extracted and saved to: {output_path.name}")
            else:
                failed += 1
                print(f"      ⚠ Failed: {result.get('error', 'Unknown error')}")
        
        # Save the completed manifest to a file
        manifest_path = self.output_dir / "_manifest.json"
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(extraction_manifest, f, indent=2)
        
        print("\n" + "="*60)
        print(f"✅ EXTRACTION COMPLETE")
        print(f"   Successful: {successful} | Failed: {failed}")
        print(f"   📖 Extraction manifest created: {manifest_path}")
        print("="*60)

if __name__ == "__main__":
    extractor = DocumentExtractor(source_dir="dataset", output_dir="output/extracted_texts")
    extractor.extract_all_documents()